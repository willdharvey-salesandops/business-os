"""
create_pptx.py

Creates the EmailShepherd qualification call deck as a PPTX and uploads
to Google Drive. On first run, creates a new file and saves its Drive ID
to drive_file_ids.json. On subsequent runs, updates the existing file.

Usage:
    python3 .claude/skills/google-slides/create_pptx.py

Environment vars required:
    GOOGLE_CREDENTIALS_PATH
    EMAILSHEPHERD_DECKS_FOLDER_ID
"""

import os
import sys
import io
import json
from pathlib import Path
from dotenv import load_dotenv

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload
from google.oauth2 import service_account

load_dotenv()

# ── Config ──────────────────────────────────────────────────────────────────

CREDENTIALS_PATH = os.getenv('GOOGLE_CREDENTIALS_PATH', '.google-credentials.json')
FOLDER_ID        = os.getenv('EMAILSHEPHERD_DECKS_FOLDER_ID')
CACHE_FILE       = Path('.claude/skills/google-slides/drive_file_ids.json')
CACHE_KEY        = ('emailshepherd', 'qualification-call-deck')
FILENAME         = 'qualification-call-deck.pptx'

# ── Colours ──────────────────────────────────────────────────────────────────

INK        = RGBColor(0x0D, 0x0D, 0x0B)
OFF_WHITE  = RGBColor(0xF8, 0xF7, 0xF4)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MUTED      = RGBColor(0x9C, 0x98, 0x90)
MID        = RGBColor(0x48, 0x45, 0x40)
BORDER     = RGBColor(0xE0, 0xDD, 0xD8)

# ── Fonts ────────────────────────────────────────────────────────────────────

DISPLAY = 'Georgia'       # approximates Fraunces
BODY    = 'Calibri'       # clean sans-serif
MONO    = 'Courier New'   # approximates DM Mono

# ── Slide dimensions: 16:9 widescreen ────────────────────────────────────────

W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_bg(slide, color):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_name=BODY, font_size=Pt(12), color=INK,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             word_wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_eyebrow(slide, text, top=Inches(1.1)):
    add_text(slide, text,
             left=Inches(1.0), top=top, width=Inches(11.33), height=Inches(0.3),
             font_name=MONO, font_size=Pt(10), color=MUTED)


def add_heading(slide, text, top=Inches(1.5), size=Pt(44), color=INK):
    add_text(slide, text,
             left=Inches(1.0), top=top, width=Inches(11.33), height=Inches(1.2),
             font_name=DISPLAY, font_size=size, color=color, bold=True)


def add_body(slide, text, top=Inches(2.9), size=Pt(14), color=MID):
    add_text(slide, text,
             left=Inches(1.0), top=top, width=Inches(11.33), height=Inches(3.5),
             font_name=BODY, font_size=size, color=color)


def add_logo(slide, dark=False):
    """Add EmailShepherd wordmark top-left."""
    color = OFF_WHITE if dark else INK
    add_text(slide, 'EmailShepherd',
             left=Inches(0.5), top=Inches(0.3), width=Inches(3), height=Inches(0.4),
             font_name=BODY, font_size=Pt(12), color=color, bold=True)


def add_label(slide, text, dark=False):
    """Add slide label top-right."""
    color = MUTED
    add_text(slide, text.upper(),
             left=Inches(9.5), top=Inches(0.3), width=Inches(3.33), height=Inches(0.4),
             font_name=MONO, font_size=Pt(9), color=color, align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, fill=WHITE):
    """Add a white card with a 2pt top accent bar."""
    # Main card
    card = add_rect(slide, left, top, width, height,
                    fill_color=fill, border_color=BORDER, border_width=Pt(0.5))
    # Top accent
    add_rect(slide, left, top, width, Pt(2), fill_color=INK)
    return card


# ── Slides ───────────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Dark title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, INK)
    add_logo(slide, dark=True)

    add_text(slide, 'DISCOVERY CALL',
             left=Inches(1.0), top=Inches(1.8), width=Inches(11.33), height=Inches(0.4),
             font_name=MONO, font_size=Pt(10), color=RGBColor(0x9C, 0x98, 0x90))

    add_text(slide, 'EmailShepherd',
             left=Inches(1.0), top=Inches(2.3), width=Inches(11.33), height=Inches(1.4),
             font_name=DISPLAY, font_size=Pt(72), color=OFF_WHITE, bold=True)

    add_text(slide, "This isn't a pitch. It's a conversation: find out if there's a fit, and if there is, agree the fastest path to knowing for certain.",
             left=Inches(1.0), top=Inches(4.1), width=Inches(8), height=Inches(1.0),
             font_name=BODY, font_size=Pt(15), color=RGBColor(0xC8, 0xC5, 0xBE))


def slide_02_process(prs):
    """9-step process overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, OFF_WHITE)
    add_logo(slide)
    add_label(slide, 'The Process')

    add_eyebrow(slide, 'How we work together', top=Inches(0.85))
    add_heading(slide, 'The path to a decision. You can opt out at any point.', top=Inches(1.2), size=Pt(32))

    add_text(slide, 'First call to onboarding in 6-10 weeks. Every step is agreed together.',
             left=Inches(1.0), top=Inches(2.2), width=Inches(11.33), height=Inches(0.35),
             font_name=BODY, font_size=Pt(11), color=MUTED)

    steps = [
        ('01', 'Qualification',    'Understand setup and fit'),
        ('02', 'Business Case',    'Build internal case for stakeholders'),
        ('03', 'Technical Demo',   'Right people, your email types'),
        ('04', 'POC Planning',     'Start date, KPIs, check-in'),
        ('05', 'Commercial Offer', 'Sent after POC planning, fully tailored'),
        ('06', 'Legal and Sign',   'POC agreement, not a commitment to buy'),
        ('07', 'POC Start',        'Access, onboarding, first emails'),
        ('08', 'POC Check-In',     'Review against agreed KPIs'),
        ('09', 'Opt Out /\nOnboarding', 'Either outcome is a good outcome'),
    ]

    tile_w = Inches(1.37)
    tile_h = Inches(1.25)
    gap    = Inches(0.08)
    start_x = Inches(0.5)
    tile_y = Inches(2.7)

    for i, (num, name, desc) in enumerate(steps):
        x = start_x + i * (tile_w + gap)
        is_active = (i == 0)
        bg = INK if is_active else WHITE
        text_c = OFF_WHITE if is_active else INK
        desc_c = RGBColor(0xC0, 0xBC, 0xB4) if is_active else MUTED

        add_card(slide, x, tile_y, tile_w, tile_h, fill=bg)

        add_text(slide, num,
                 left=x, top=tile_y + Inches(0.12), width=tile_w, height=Inches(0.25),
                 font_name=MONO, font_size=Pt(8), color=desc_c, align=PP_ALIGN.CENTER)

        add_text(slide, name,
                 left=x + Inches(0.06), top=tile_y + Inches(0.4), width=tile_w - Inches(0.12), height=Inches(0.5),
                 font_name=BODY, font_size=Pt(10), color=text_c, bold=True, align=PP_ALIGN.CENTER)

        desc_y = tile_y + tile_h + Inches(0.05)
        add_text(slide, desc,
                 left=x, top=desc_y, width=tile_w, height=Inches(0.55),
                 font_name=BODY, font_size=Pt(8), color=MUTED, align=PP_ALIGN.CENTER)

    # Callout
    callout_y = Inches(5.85)
    add_card(slide, Inches(0.5), callout_y, Inches(12.33), Inches(0.8))
    add_text(slide, 'Signing the POC agreement is not a commitment to buy.  It is a mutual framework for testing fairly, with agreed start dates, KPIs, and a check-in. At the end, you decide: opt out, or move to full onboarding.',
             left=Inches(0.65), top=callout_y + Inches(0.15), width=Inches(12.0), height=Inches(0.6),
             font_name=BODY, font_size=Pt(10), color=MID)


def slide_03_what_we_do(prs):
    """What EmailShepherd does — two product pillars."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, OFF_WHITE)
    add_logo(slide)
    add_label(slide, 'What We Do')

    add_eyebrow(slide, 'The product')
    add_heading(slide, 'Email Shepherd', top=Inches(1.5), size=Pt(48))

    add_text(slide, 'An email design system that sits between your brand guidelines and your ESP.',
             left=Inches(1.0), top=Inches(2.5), width=Inches(11.33), height=Inches(0.5),
             font_name=BODY, font_size=Pt(15), color=MID)

    cards = [
        ('Email Design System',
         'Modular components, full HTML control. One source of truth for every email, no code knowledge required for marketers.'),
        ('EmailShepherd AI',
         'Generate on-brand emails from a prompt. Translate instantly. Automated QA and review, without the hours.'),
    ]

    card_w = Inches(5.5)
    card_h = Inches(2.2)
    y = Inches(3.3)
    for i, (title, body) in enumerate(cards):
        x = Inches(1.0) + i * (card_w + Inches(0.33))
        add_card(slide, x, y, card_w, card_h)
        add_text(slide, title,
                 left=x + Inches(0.25), top=y + Inches(0.3), width=card_w - Inches(0.5), height=Inches(0.4),
                 font_name=BODY, font_size=Pt(13), color=INK, bold=True)
        add_text(slide, body,
                 left=x + Inches(0.25), top=y + Inches(0.8), width=card_w - Inches(0.5), height=Inches(1.2),
                 font_name=BODY, font_size=Pt(11), color=MID)


def slide_04_who_we_work_with(prs):
    """ICP slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, OFF_WHITE)
    add_logo(slide)
    add_label(slide, 'Who We Work With')

    add_eyebrow(slide, 'Ideal fit')
    add_heading(slide, "Before I ask about your setup.\nHere's who typically comes to us.", top=Inches(1.5), size=Pt(36))

    cards = [
        ('On a platform that has been acquired',
         'Taxi for Email was acquired twice. Now it\'s Bird. Platform uncertainty is a real operational and commercial risk.'),
        ('Scaling across teams or markets',
         'Multiple brands, regions, or languages. Without a single source of truth, brand consistency breaks. And legally, that matters.'),
    ]

    card_w = Inches(5.5)
    card_h = Inches(2.2)
    y = Inches(3.7)
    for i, (title, body) in enumerate(cards):
        x = Inches(1.0) + i * (card_w + Inches(0.33))
        add_card(slide, x, y, card_w, card_h)
        add_text(slide, title,
                 left=x + Inches(0.25), top=y + Inches(0.3), width=card_w - Inches(0.5), height=Inches(0.4),
                 font_name=BODY, font_size=Pt(13), color=INK, bold=True)
        add_text(slide, body,
                 left=x + Inches(0.25), top=y + Inches(0.8), width=card_w - Inches(0.5), height=Inches(1.2),
                 font_name=BODY, font_size=Pt(11), color=MID)


def slide_05_current_setup(prs):
    """Discovery questions — current setup."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, OFF_WHITE)
    add_logo(slide)
    add_label(slide, 'Your Setup')

    add_eyebrow(slide, 'Discovery')
    add_heading(slide, 'Walk me through how you build an email today.', top=Inches(1.5), size=Pt(40))

    questions = [
        ('Process',      'Step by step, from brief to send. Who does what?'),
        ('Team',         'How many people are involved? Where does it get stuck?'),
        ('Tools',        'What are you using today? What does it do well or badly?'),
        ('Volume',       'How many emails a month? How many hours per email, roughly?'),
        ('Why now',      'Have you tried to change this before? What is different about now?'),
    ]

    q_y = Inches(2.6)
    q_h = Inches(0.6)
    for i, (label, question) in enumerate(questions):
        y = q_y + i * (q_h + Inches(0.08))
        add_rect(slide, Inches(1.0), y, Inches(11.33), q_h, fill_color=WHITE,
                 border_color=BORDER, border_width=Pt(0.5))
        add_text(slide, label,
                 left=Inches(1.15), top=y + Inches(0.1), width=Inches(1.2), height=Inches(0.4),
                 font_name=MONO, font_size=Pt(9), color=MUTED)
        add_text(slide, question,
                 left=Inches(2.5), top=y + Inches(0.1), width=Inches(9.5), height=Inches(0.4),
                 font_name=BODY, font_size=Pt(12), color=INK, bold=True)


def slide_06_common_friction(prs):
    """Common pains slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, OFF_WHITE)
    add_logo(slide)
    add_label(slide, 'Common Friction')

    add_eyebrow(slide, 'Where teams usually feel the pain')
    add_heading(slide, 'Does any of this sound familiar?', top=Inches(1.5), size=Pt(40))

    pains = [
        ('1', 'Hours lost to manual QA: checking brand, accessibility, and rendering on every email'),
        ('2', 'Designer bottleneck: marketers waiting on dev or design for every small change'),
        ('3', 'Brand inconsistency: different teams or regions not looking like the same company'),
        ('4', 'ESP limitations: working around what the tool cannot do instead of what the brand needs'),
        ('5', 'Platform risk: depending on a tool that has been acquired, roadmap uncertain'),
    ]

    p_y = Inches(2.6)
    p_h = Inches(0.65)
    for i, (num, pain) in enumerate(pains):
        y = p_y + i * (p_h + Inches(0.05))
        add_rect(slide, Inches(1.0), y, Inches(11.33), p_h, fill_color=WHITE,
                 border_color=BORDER, border_width=Pt(0.5))
        add_text(slide, num,
                 left=Inches(1.15), top=y + Inches(0.12), width=Inches(0.5), height=Inches(0.4),
                 font_name=BODY, font_size=Pt(16), color=INK, bold=True)
        add_text(slide, pain,
                 left=Inches(1.8), top=y + Inches(0.12), width=Inches(10.3), height=Inches(0.4),
                 font_name=BODY, font_size=Pt(12), color=INK)


def slide_07_bbc(prs):
    """BBC case study."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, INK)
    add_logo(slide, dark=True)

    add_text(slide, 'CASE STUDY',
             left=Inches(1.0), top=Inches(0.85), width=Inches(11.33), height=Inches(0.3),
             font_name=MONO, font_size=Pt(10), color=MUTED)

    add_text(slide, 'BBC',
             left=Inches(1.0), top=Inches(1.3), width=Inches(6), height=Inches(0.8),
             font_name=DISPLAY, font_size=Pt(64), color=OFF_WHITE, bold=True)

    add_text(slide, '6 brands',
             left=Inches(1.0), top=Inches(2.2), width=Inches(6), height=Inches(0.8),
             font_name=DISPLAY, font_size=Pt(56), color=OFF_WHITE, bold=True)

    add_text(slide, 'Rolled out across News, Sport, iPlayer and more. In one month.',
             left=Inches(1.0), top=Inches(3.1), width=Inches(6), height=Inches(0.6),
             font_name=BODY, font_size=Pt(13), color=RGBColor(0x9C, 0x98, 0x90))

    add_text(slide, '"Insanely fast and easy to put together."',
             left=Inches(1.0), top=Inches(4.2), width=Inches(9), height=Inches(0.7),
             font_name=DISPLAY, font_size=Pt(22), color=OFF_WHITE, italic=True)

    add_text(slide, 'First feedback from a BBC user',
             left=Inches(1.0), top=Inches(4.9), width=Inches(6), height=Inches(0.4),
             font_name=BODY, font_size=Pt(11), color=MUTED)

    stats = [
        ('17', 'users onboarded'),
        ('0', 'broken emails'),
        ('1 month', 'to roll out 6 brands'),
    ]
    for i, (val, label) in enumerate(stats):
        x = Inches(8.5) + i * Inches(0.0)
        stat_y = Inches(2.3 + i * 1.2)
        add_text(slide, val,
                 left=Inches(8.5), top=stat_y, width=Inches(4.5), height=Inches(0.55),
                 font_name=DISPLAY, font_size=Pt(36), color=OFF_WHITE, bold=True)
        add_text(slide, label,
                 left=Inches(8.5), top=stat_y + Inches(0.55), width=Inches(4.5), height=Inches(0.3),
                 font_name=BODY, font_size=Pt(11), color=MUTED)


def slide_08_retailer(prs):
    """Major retailer case study."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, INK)
    add_logo(slide, dark=True)

    add_text(slide, 'CASE STUDY',
             left=Inches(1.0), top=Inches(0.85), width=Inches(11.33), height=Inches(0.3),
             font_name=MONO, font_size=Pt(10), color=MUTED)

    add_text(slide, 'Major Online Retailer',
             left=Inches(1.0), top=Inches(1.3), width=Inches(9), height=Inches(0.8),
             font_name=DISPLAY, font_size=Pt(48), color=OFF_WHITE, bold=True)

    add_text(slide, '2 weeks to 1 day',
             left=Inches(1.0), top=Inches(2.2), width=Inches(7), height=Inches(0.8),
             font_name=DISPLAY, font_size=Pt(52), color=OFF_WHITE, bold=True)

    before_after = [
        ('Before', '2-week agency turnaround. Generic selling. Impossible to react to market.'),
        ('After',  '24-hour turnaround. Real-time data feeds. Agency removed. Revenue grew immediately.'),
    ]
    for i, (label, text) in enumerate(before_after):
        y = Inches(3.2 + i * 1.1)
        add_text(slide, label,
                 left=Inches(1.0), top=y, width=Inches(1.2), height=Inches(0.35),
                 font_name=MONO, font_size=Pt(9), color=MUTED)
        add_text(slide, text,
                 left=Inches(2.3), top=y, width=Inches(10.7), height=Inches(0.7),
                 font_name=BODY, font_size=Pt(13), color=RGBColor(0xC8, 0xC5, 0xBE))

    add_text(slide, '"They sold more straight away."',
             left=Inches(1.0), top=Inches(5.5), width=Inches(9), height=Inches(0.7),
             font_name=DISPLAY, font_size=Pt(22), color=OFF_WHITE, italic=True)


def slide_09_questions(prs):
    """Four qualification questions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, OFF_WHITE)
    add_logo(slide)
    add_label(slide, 'Qualify')

    add_eyebrow(slide, 'Before we talk next steps')
    add_heading(slide, 'Four things I want to understand.', top=Inches(1.5), size=Pt(40))

    questions = [
        ('Decision Process', 'How does a decision like this typically get made, and who signs off?'),
        ('Urgency',          'What is prompting this now? What changes if you do nothing in the next 12 months?'),
        ('Stakeholders',     'Who else would be involved in evaluating and approving a change like this?'),
        ('Success',          'If we ran a POC, what would it need to show for you to feel confident?'),
    ]

    q_y = Inches(2.6)
    q_h = Inches(0.85)
    q_gap = Inches(0.12)
    for i, (label, question) in enumerate(questions):
        y = q_y + i * (q_h + q_gap)
        add_card(slide, Inches(1.0), y, Inches(11.33), q_h)
        add_text(slide, label,
                 left=Inches(1.2), top=y + Inches(0.1), width=Inches(2.5), height=Inches(0.3),
                 font_name=MONO, font_size=Pt(9), color=MUTED)
        add_text(slide, question,
                 left=Inches(1.2), top=y + Inches(0.42), width=Inches(11.0), height=Inches(0.35),
                 font_name=BODY, font_size=Pt(13), color=INK, bold=True)


def slide_10_next_steps(prs):
    """Dark close slide with next steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, INK)
    add_logo(slide, dark=True)

    add_text(slide, 'WHAT HAPPENS NEXT',
             left=Inches(1.0), top=Inches(0.85), width=Inches(11.33), height=Inches(0.3),
             font_name=MONO, font_size=Pt(10), color=MUTED)

    add_heading(slide, "Let's agree on the next step.", top=Inches(1.3), size=Pt(44), color=OFF_WHITE)

    steps = [
        ('done',  'Today: Discovery',   'Understand your current setup and challenges',                    'Done'),
        ('now',   'Business Case',      'Build internal case, map stakeholders, sense check budget',       '~Week 2'),
        ('future','Technical Demo',     'Right people in the room, your emails, your questions answered',  '~Week 3'),
        ('future','POC Planning to Onboarding', 'Agreed KPIs, POC agreement (not a commitment to buy), 4-6 week POC, you decide', 'Weeks 4-10'),
    ]

    s_y = Inches(2.5)
    s_h = Inches(0.8)
    for i, (state, title, desc, time) in enumerate(steps):
        y = s_y + i * (s_h + Inches(0.1))
        bg = RGBColor(0x2A, 0x2A, 0x27) if state == 'now' else RGBColor(0x1A, 0x1A, 0x17)
        border_c = RGBColor(0x3A, 0x3A, 0x37)
        add_rect(slide, Inches(1.0), y, Inches(11.33), s_h,
                 fill_color=bg, border_color=border_c, border_width=Pt(0.5))

        num_text = 'v' if state == 'done' else str(i + 1)
        num_color = MUTED if state == 'done' else (OFF_WHITE if state == 'now' else MUTED)
        add_text(slide, num_text,
                 left=Inches(1.15), top=y + Inches(0.15), width=Inches(0.5), height=Inches(0.5),
                 font_name=BODY, font_size=Pt(14), color=num_color, bold=True)

        title_color = MUTED if state == 'done' else OFF_WHITE
        add_text(slide, title,
                 left=Inches(1.8), top=y + Inches(0.08), width=Inches(7.5), height=Inches(0.35),
                 font_name=BODY, font_size=Pt(13), color=title_color, bold=True)
        add_text(slide, desc,
                 left=Inches(1.8), top=y + Inches(0.42), width=Inches(7.5), height=Inches(0.3),
                 font_name=BODY, font_size=Pt(10), color=MUTED)
        add_text(slide, time,
                 left=Inches(10.5), top=y + Inches(0.2), width=Inches(1.6), height=Inches(0.4),
                 font_name=MONO, font_size=Pt(10), color=MUTED, align=PP_ALIGN.RIGHT)

    add_text(slide, 'What does your diary look like next week?',
             left=Inches(1.0), top=Inches(6.7), width=Inches(11.33), height=Inches(0.5),
             font_name=DISPLAY, font_size=Pt(18), color=RGBColor(0x60, 0x5D, 0x58),
             italic=True, align=PP_ALIGN.CENTER)


# ── Build presentation ────────────────────────────────────────────────────────

def build_pptx():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_01_title(prs)
    slide_02_process(prs)
    slide_03_what_we_do(prs)
    slide_04_who_we_work_with(prs)
    slide_05_current_setup(prs)
    slide_06_common_friction(prs)
    slide_07_bbc(prs)
    slide_08_retailer(prs)
    slide_09_questions(prs)
    slide_10_next_steps(prs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── Drive helpers ─────────────────────────────────────────────────────────────

def get_drive_service():
    scopes = ['https://www.googleapis.com/auth/drive']
    creds = service_account.Credentials.from_service_account_file(
        CREDENTIALS_PATH, scopes=scopes
    )
    return build('drive', 'v3', credentials=creds)


def load_cache():
    if CACHE_FILE.exists():
        return json.loads(CACHE_FILE.read_text())
    return {}


def save_cache(cache):
    CACHE_FILE.parent.mkdir(parents=True, exist_ok=True)
    CACHE_FILE.write_text(json.dumps(cache, indent=2))


def upload_pptx(service, buf, folder_id, existing_id=None):
    mime = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    media = MediaIoBaseUpload(buf, mimetype=mime, resumable=False)

    if existing_id:
        file = service.files().update(
            fileId=existing_id,
            media_body=media,
            fields='id, webViewLink',
            supportsAllDrives=True
        ).execute()
        print(f"Updated existing file: {existing_id}")
    else:
        metadata = {
            'name': FILENAME,
            'parents': [folder_id],
        }
        file = service.files().create(
            body=metadata,
            media_body=media,
            fields='id, webViewLink',
            supportsAllDrives=True
        ).execute()
        print(f"Created new file: {file['id']}")

    return file


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    if not FOLDER_ID:
        print("ERROR: EMAILSHEPHERD_DECKS_FOLDER_ID not set in .env")
        print("Run setup_es_folders.py first and add the folder IDs to .env")
        sys.exit(1)

    print("Building PPTX...")
    buf = build_pptx()

    print("Connecting to Drive...")
    service = get_drive_service()

    cache = load_cache()
    client_cache = cache.get(CACHE_KEY[0], {})
    existing_id = client_cache.get(CACHE_KEY[1])

    print("Uploading to Drive...")
    file = upload_pptx(service, buf, FOLDER_ID, existing_id)

    # Update cache
    if CACHE_KEY[0] not in cache:
        cache[CACHE_KEY[0]] = {}
    cache[CACHE_KEY[0]][CACHE_KEY[1]] = file['id']
    save_cache(cache)

    link = file.get('webViewLink', f"https://drive.google.com/file/d/{file['id']}/view")
    print(f"\nDone. Drive link:\n{link}")


if __name__ == '__main__':
    main()
